"""JSON → .pptx 课件渲染器（真实课件视觉体系）。

设计语言对齐优质小学课件成品：
- 暖白底 + 近黑粗体大标题 + 标题下强调色短横（无通栏色带）
- 白色圆角卡 + 柔和外阴影 + 顶部彩带
- 关键词彩色高亮（emphasize 逐 run 上色，蓝/绿/红/紫轮换）
- 序号药丸徽章、点线引用框、树形板书（节点框 + 连线）
- 整页页脚（左课程名 · 右页码），右上角课时/栏目药丸标签
- 注音行 / 田字格用 Pillow 预渲染 PNG

退出码：0 成功 / 1 异常 / 2 前置缺失
"""
from __future__ import annotations
import io
import sys
import os
from pathlib import Path

# Windows 控制台默认 GBK，输出 ✓/✗ 等 Unicode 会崩；强制 UTF-8
try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

from common import design_tokens as T
from common import fonts, pinyin as py
from common.schema import stage_short

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE
from pptx.oxml.ns import qn
from lxml import etree

# ---- 颜色/字体辅助 ----
def _rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def _round_rect(shape, radius=0.10):
    """设置圆角矩形圆角比例。"""
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass


def _soft_shadow(shape, blur=0.055, dist=0.035, direction=5400000,
                 color=None, alpha=42):
    """给 autoshape 注入柔和外阴影（python-pptx 无原生 API，手写 XML）。"""
    color = color or T.PAL.FOOT
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    xml = (
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{Emu(int(blur * 914400))}" '
        f'dist="{Emu(int(dist * 914400))}" dir="{direction}" rotWithShape="0">'
        f'<a:srgbClr val="{color}"><a:alpha val="{int(alpha * 1000)}"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    )
    try:
        spPr.append(etree.fromstring(xml))
    except Exception:
        pass


def _no_border(shape):
    shape.line.fill.background()


def _add_textbox(slide, left, top, width, height, text="", *,
                 font=fonts.HEI, size=24, color=None, bold=False,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 runs=None, line_spacing=None):
    """添加文本框。runs=[(text, {size,color,bold,font}), ...] 优先于 text，
    用于同段落内多色关键词高亮。"""
    color = color or T.PAL.TEXT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if runs:
        for rtext, rstyle in runs:
            run = p.add_run()
            run.text = rtext
            run.font.name = rstyle.get("font", font)
            run.font.size = Pt(rstyle.get("size", size))
            run.font.color.rgb = _rgb(rstyle.get("color", color))
            run.font.bold = rstyle.get("bold", bold)
    else:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(color)
        run.font.bold = bold
    return tb


def _bg(slide, color=None):
    """整页暖白背景矩形（置底）。"""
    color = color or T.PAL.BG
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, T.SLIDE_W, T.SLIDE_H)
    shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(color)
    _no_border(shp)
    shp.shadow.inherit = False
    sp = shp._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shp


def _accent_underline(slide, x, y, w=0.85, h=0.11, color=None):
    """标题下方强调色短横。"""
    color = color or T.PAL.ACCENT
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 int(x), int(y), inch(w), inch(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(color)
    _no_border(bar); bar.shadow.inherit = False
    _round_rect(bar, 0.5)
    return bar


def inch(v):
    return T.inch(v)


# ---- 页眉 / 页脚 / 标签 ----
KIND_LABEL = {
    "cover": "封面", "toc": "目录", "objectives": "学习目标", "intro": "情境导入",
    "reading": "朗读感知", "word-cards": "识字卡片", "words": "词语天地",
    "comprehension": "精读品析", "summary": "课堂小结", "homework": "作业布置",
    "review": "复习导入", "writing": "写法总结", "extend": "拓展延伸",
    "board": "板书设计", "practice": "课堂练习", "end": "结束",
    "revision": "写字指导", "poem": "古诗", "discussion": "合作探究",
}


def _header(slide, sj, stage):
    """页标题体系：左上大粗标题 + 下方橙色短横 + 右上药丸标签。"""
    _accent_underline(slide, inch(T.L.MARGIN_X), inch(1.12), w=0.95)
    _add_textbox(slide, inch(T.L.MARGIN_X), inch(0.42), inch(9.0), inch(0.72),
                 sj.get("title", ""), font=fonts.HEI,
                 size=T.font_for(stage, "slide_title"),
                 color=T.PAL.TITLE_TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)


def _pill(slide, x_right, y, text, *, fill=None, color=None, size=13, bold=True,
          pad=0.16, h=0.34):
    """右对齐药丸标签，返回其左边界 EMU。"""
    fill = fill or T.PAL.ACCENT
    color = color or T.PAL.TITLE_TEXT
    w = 0.20 + len(text) * (size * 0.0135) + pad * 2
    left = x_right - inch(w)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  int(left), int(y), inch(w), inch(h))
    _round_rect(pill, 0.5)
    if fill == "line":
        pill.fill.background()
        pill.line.color.rgb = _rgb(color); pill.line.width = Pt(1.4)
    else:
        pill.fill.solid(); pill.fill.fore_color.rgb = _rgb(fill)
        _no_border(pill)
    pill.shadow.inherit = False
    _add_textbox(slide, int(left), int(y), inch(w), inch(h), text,
                 font=fonts.HEI, size=size, color=color, bold=bold,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return left


def _header_tags(slide, sj, meta, total_periods):
    """右上角：栏目名（描边药丸）+ 课时（实心药丸）。"""
    x = inch(13.333 - 0.55)
    y = inch(0.62)
    per = sj.get("period", 1)
    if total_periods > 1:
        left = _pill(slide, x, y, f"第{per}课时", fill=T.PAL.ACCENT, size=13)
    else:
        left = x - inch(0.2)
    kind = sj.get("kind", "")
    label = KIND_LABEL.get(kind)
    if label:
        _pill(slide, left - inch(0.12), y, label, fill="line",
              color=T.PAL.ACCENT, size=13)


def _footer(slide, meta, page_no, total):
    """页脚：左课程名 · 右页码，细灰。"""
    y = inch(7.16)
    _add_textbox(slide, inch(T.L.MARGIN_X), y, inch(6.0), inch(0.3),
                 f"{meta.get('title','')} · {meta.get('textbook','')}",
                 font=fonts.HEI, size=11, color=T.PAL.FOOT,
                 anchor=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, inch(11.2), y, inch(1.6), inch(0.3),
                 f"{page_no:02d} / {total:02d}", font=fonts.HEI, size=11,
                 color=T.PAL.FOOT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ---- Pillow 预渲染：注音行 / 田字格 ----
def _render_ruby_png(text, pinyin_str, stage, big=False):
    """整行注音 → PNG bytes。Pillow 画拼音(小)+汉字(大)两行。"""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return None
    pairs = py.split_syllables(text, pinyin_str)
    char_sz = T.font_for(stage, "bigchar" if big else "body") + (8 if big else 0)
    py_sz = T.font_for(stage, "pinyin")
    pad = 10
    cell_w = char_sz + 18
    w = max(cell_w * len(pairs) + pad * 2, 200)
    h = py_sz + char_sz + pad * 3
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    cf = _load_font(fonts.HEI, char_sz)
    pf = _load_font(fonts.MONO, py_sz)
    x = pad
    for c, s, tone in pairs:
        color = "#" + py.tone_color(tone)
        bw = draw.textlength(s, font=pf) if pf else len(s) * py_sz
        draw.text((x + (cell_w - bw) / 2, pad), s, fill=color, font=pf)
        cw = draw.textlength(c, font=cf) if cf else char_sz
        draw.text((x + (cell_w - cw) / 2, pad + py_sz + 6), c,
                  fill="#" + T.PAL.TITLE_TEXT, font=cf)
        x += cell_w
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def _render_tianzi_png(char, stage, tone_color=None):
    """田字格 + 大字 → PNG bytes（透明底，虚线米字格）。"""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return None
    sz = T.font_for(stage, "bigchar")
    size = int(sz * 1.7) + 44
    img = Image.new("RGBA", (size, size), (255, 255, 255, 255))
    draw = ImageDraw.Draw(img)
    grid = "#" + T.PAL.ACCENT
    m = 10
    # 外框
    draw.rectangle([m, m, size - m, size - m], outline="#D9C6A5", width=3)
    # 十字虚线
    cx = cy = size // 2
    dash = 10
    for y in range(m + 6, size - m - 6, dash * 2):
        draw.line([(cx, y), (cx, min(y + dash, size - m - 6))], fill=grid, width=2)
    for x in range(m + 6, size - m - 6, dash * 2):
        draw.line([(x, cy), (min(x + dash, size - m - 6), cy)], fill=grid, width=2)
    # 对角虚线
    for t in range(0, size - 2 * m - 12, dash * 2):
        draw.line([(m + 6 + t, m + 6 + t),
                   (min(m + 6 + t + dash, size - m - 6), min(m + 6 + t + dash, size - m - 6))],
                  fill="#EAD9BE", width=1)
        draw.line([(size - m - 6 - t, m + 6 + t),
                   (max(size - m - 6 - t - dash, m + 6), min(m + 6 + t + dash, size - m - 6))],
                  fill="#EAD9BE", width=1)
    f = _load_font(fonts.KAI, int(sz * 1.25))
    color = "#" + (tone_color or T.PAL.TITLE_TEXT)
    if f:
        bb = draw.textbbox((0, 0), char, font=f)
        cw, ch = bb[2] - bb[0], bb[3] - bb[1]
        draw.text(((size - cw) / 2 - bb[0], (size - ch) / 2 - bb[1]), char, fill=color, font=f)
    else:
        draw.text((size * 0.28, size * 0.22), char, fill=color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


_FONT_CACHE = {}
def _load_font(name, size):
    """尝试加载系统字体，失败返回 None（退化为默认）。"""
    key = (name, size)
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    try:
        from PIL import ImageFont
    except ImportError:
        _FONT_CACHE[key] = None
        return None
    candidates = []
    if os.name == "nt":
        windir = os.environ.get("WINDIR", r"C:\Windows")
        name_map = {"微软雅黑": "msyh.ttc", "Microsoft YaHei": "msyh.ttc",
                    "宋体": "simsun.ttc", "SimSun": "simsun.ttc",
                    "楷体": "simkai.ttf", "KaiTi": "simkai.ttf",
                    "黑体": "simhei.ttf", "SimHei": "simhei.ttf",
                    "Consolas": "consola.ttf"}
        fn = name_map.get(name)
        if fn:
            candidates.append(os.path.join(windir, "Fonts", fn))
        candidates += [os.path.join(windir, "Fonts", "msyh.ttc"),
                       os.path.join(windir, "Fonts", "simsun.ttc")]
    for path in candidates:
        if os.path.exists(path):
            try:
                f = ImageFont.truetype(path, size)
                _FONT_CACHE[key] = f
                return f
            except Exception:
                continue
    _FONT_CACHE[key] = None
    return None


# ---- 高亮 run 拆分 ----
def _hl_runs(content, emphasize, base_size, base_color=None):
    """把 content 按 emphasize[{start,end,color?,bold?}] 切成多色 runs。"""
    base_color = base_color or T.PAL.TEXT
    if not content:
        return [(content or "", {})]
    marks = []
    for i, em in enumerate(emphasize or []):
        try:
            s, e = int(em.get("start", 0)), int(em.get("end", 0))
        except (TypeError, ValueError):
            continue
        if e <= s:
            continue
        color = em.get("color") or T.PAL.HIGHLIGHTS[i % len(T.PAL.HIGHLIGHTS)]
        marks.append((max(0, min(s, len(content))),
                      max(0, min(e, len(content))), color, em.get("bold", True)))
    marks.sort()
    runs, pos = [], 0
    for s, e, color, bold in marks:
        if s < pos:
            continue
        if s > pos:
            runs.append((content[pos:s], {"size": base_size, "color": base_color}))
        runs.append((content[s:e], {"size": base_size, "color": color, "bold": bold}))
        pos = e
    if pos < len(content):
        runs.append((content[pos:], {"size": base_size, "color": base_color}))
    return runs or [(content, {})]


# ---- 布局分类 ----
def classify_slide(slide_json) -> str:
    kind = slide_json.get("kind", "")
    elems = slide_json.get("elements", [])
    types = [e.get("type") for e in elems]
    if kind in ("cover", "end"):
        return "cover"
    if "word-card" in types:
        return "cards"
    if "revision" in types:
        return "revision"
    if "poem" in types or "ruby-line" in types:
        return "reading"
    if "board" in types:
        return "board"
    if "table" in types:
        return "table"
    return "text"


# ---- 页面渲染 ----
def _render_cover(slide, sj, meta, stage, page_no, total):
    """封面：装饰圆 + 大字标题 + 短横 + 信息行。"""
    _bg(slide)
    # 顶部装饰：细横条 + 三圆点
    _accent_underline(slide, inch(5.9), inch(1.15), w=1.5, h=0.06)
    for i, (cx, col) in enumerate([(5.55, T.PAL.ACCENT), (6.60, T.PAL.ACCENT2), (7.65, T.PAL.ACCENT3)]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx), inch(1.5), inch(0.16), inch(0.16))
        dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(col)
        _no_border(dot); dot.shadow.inherit = False
    title = sj.get("title") or meta["title"]
    # 标题下大短横
    _accent_underline(slide, inch(5.42), inch(4.30), w=2.5, h=0.12)
    _add_textbox(slide, inch(1.2), inch(2.35), inch(10.93), inch(1.9),
                 title, font=fonts.HEI, size=T.font_for(stage, "cover_title"),
                 color=T.PAL.TITLE_TEXT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    sub = f"{meta.get('textbook','')}　·　{meta.get('lessonType','')}课"
    _add_textbox(slide, inch(1.2), inch(4.62), inch(10.93), inch(0.7),
                 sub, font=fonts.HEI, size=22, color=T.PAL.TEXT_LIGHT,
                 align=PP_ALIGN.CENTER)
    per = sj.get("period", 1)
    if meta.get("periods", 1) > 1:
        _pill(slide, inch(6.9), inch(5.5), f"第 {per} 课时",
              fill=T.PAL.ACCENT, size=15)
    _footer(slide, meta, page_no, total)


def _render_end(slide, sj, meta, stage, page_no, total):
    _bg(slide)
    title = sj.get("title") or "本课结束"
    _accent_underline(slide, inch(5.92), inch(3.9), w=1.5, h=0.10)
    _add_textbox(slide, inch(1.2), inch(2.6), inch(10.93), inch(1.3),
                 title, font=fonts.HEI, size=T.font_for(stage, "cover_title") - 10,
                 color=T.PAL.ACCENT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    top = inch(T.L.CONTENT_TOP)
    for el in sj.get("elements", []):
        if el.get("type") in ("paragraph", "quote", "list"):
            _place_element(slide, el, top, stage, 0)
            break
    _footer(slide, meta, page_no, total)


def _render_content(slide, sj, meta, stage):
    _bg(slide)
    _header(slide, sj, stage)
    top = inch(T.L.CONTENT_TOP)
    for i, el in enumerate(sj.get("elements", [])):
        if top > inch(T.L.MAX_Y):
            break
        top = _place_element(slide, el, top, stage, i)


# ---- 通用元素放置 ----
def _card_panel(slide, x, y, w, h, *, fill=None, edge_color=None, edge_w=None):
    """白色圆角卡 + 柔影。"""
    fill = fill or T.PAL.BG_CARD
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    _round_rect(card, 0.10)
    card.fill.solid(); card.fill.fore_color.rgb = _rgb(fill)
    if edge_color:
        card.line.color.rgb = _rgb(edge_color)
        card.line.width = Pt(edge_w or 1.2)
    else:
        _no_border(card)
    card.shadow.inherit = False
    _soft_shadow(card)
    return card


def _place_element(slide, el, top, stage, idx=0) -> int:
    t = el.get("type")
    left = inch(T.L.MARGIN_X)
    width = inch(13.333 - 2 * T.L.MARGIN_X)
    gap = inch(T.L.GAP)

    if t == "heading":
        sz = T.font_for(stage, el.get("size", "h2"))
        # 小竖条 + 栏内标题
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     int(left), int(top + inch(0.06)), inch(0.09), inch(sz / 72 * 1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(T.PAL.ACCENT)
        _no_border(bar); bar.shadow.inherit = False
        _add_textbox(slide, left + inch(0.24), top, width, inch(0.62),
                     el.get("content", ""), font=fonts.HEI, size=sz,
                     color=T.PAL.TITLE_TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        return top + inch(0.72) + gap

    if t == "paragraph":
        sz = T.font_for(stage, "body")
        content = el.get("content", "")
        runs = _hl_runs(content, el.get("emphasize"), sz)
        # 高度估算：按可视字符宽度
        width_in = 13.333 - 2 * T.L.MARGIN_X
        chars_per_line = max(1, int(width_in * 72 / (sz * 1.02)))
        lines = max(1, -(-len(content) // chars_per_line))
        h = inch(0.42 * lines + 0.24)
        _add_textbox(slide, left, top, width, h, runs=runs,
                     font=fonts.HEI, size=sz, color=T.PAL.TEXT,
                     anchor=MSO_ANCHOR.TOP, line_spacing=T.L.LINE_SP)
        return top + h + gap

    if t == "list":
        return _place_list(slide, el, top, stage) + gap

    if t == "quote":
        return _place_quote(slide, el, top, stage) + gap

    if t == "table":
        return _place_table(slide, el, top, stage) + gap

    if t == "ruby-line":
        return _place_ruby_line(slide, el, top, stage) + gap

    if t == "poem":
        return _place_poem(slide, el, top, stage) + gap

    if t == "word-card":
        return _place_word_cards(slide, el, top, stage) + gap

    if t == "revision":
        return _place_revision(slide, el, top, stage) + gap

    if t == "board":
        return _place_board(slide, el, top, stage) + gap

    if t == "discussion":
        return _place_discussion(slide, el, top, stage) + gap

    if t == "evaluation":
        return _place_evaluation(slide, el, top, stage) + gap

    if t == "strokes":
        return _place_strokes(slide, el, top, stage) + gap

    if t == "divider":
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, T.inch(4), T.inch(3.3),
                                     T.inch(5.3), T.inch(0.04))
        bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(T.PAL.DIVIDER)
        _no_border(bar); bar.shadow.inherit = False
        return top + inch(1.0)

    if t == "note":
        _add_textbox(slide, left, top, width, inch(0.4),
                     "◈ 教师备注：" + el.get("content", ""),
                     font=fonts.HEI, size=13, color=T.PAL.TEXT_LIGHT)
        return top + inch(0.42) + gap

    if t == "image":
        # 无网络/本地素材时留柔和占位面板
        w_in = 13.333 - 2 * T.L.MARGIN_X
        h = inch(el.get("height", 1.6))
        _card_panel(slide, left, top, inch(w_in), h, fill="FFF3E2")
        _add_textbox(slide, left, top, inch(w_in), h,
                     "🖼  " + (el.get("caption") or "插图（请在 PPT 中替换为实拍图）"),
                     font=fonts.HEI, size=14, color=T.PAL.TEXT_LIGHT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return top + h + gap

    return top + gap


def _place_list(slide, el, top, stage):
    """目录/目标列表：橙色数字徽章或彩色圆点 + 文本，整组白卡。"""
    items = el.get("items", [])
    if not items:
        return top
    ordered = el.get("ordered", False)
    sz = T.font_for(stage, "list")
    row_h = 0.62 if len(items) <= 6 else 0.5
    h = inch(row_h * len(items) + 0.4)
    left = inch(T.L.MARGIN_X)
    width = inch(13.333 - 2 * T.L.MARGIN_X)
    _card_panel(slide, left, top, width, h)
    y = top + inch(0.2)
    for i, item in enumerate(items):
        col = T.PAL.HIGHLIGHTS[i % len(T.PAL.HIGHLIGHTS)] if ordered else T.PAL.ACCENT
        ix = left + inch(0.35)
        if ordered:
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(ix), int(y + inch((row_h - 0.42) / 2)),
                                           inch(0.42), inch(0.42))
            badge.fill.solid(); badge.fill.fore_color.rgb = _rgb(T.PAL.ACCENT)
            _no_border(badge); badge.shadow.inherit = False
            _add_textbox(slide, int(ix), int(y + inch((row_h - 0.42) / 2)), inch(0.42), inch(0.42),
                         str(i + 1), font=fonts.HEI, size=sz - 4,
                         color=T.PAL.TITLE_TEXT, bold=True,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            tx = ix + inch(0.62)
        else:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(ix + inch(0.1)),
                                         int(y + inch(row_h / 2 - 0.055)), inch(0.11), inch(0.11))
            dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(T.PAL.ACCENT)
            _no_border(dot); dot.shadow.inherit = False
            tx = ix + inch(0.38)
        runs = _hl_runs(str(item), el.get("emphasize"), sz) if isinstance(item, str) else [(str(item), {})]
        _add_textbox(slide, tx, int(y), width - (tx - left) - inch(0.3), inch(row_h),
                     runs=runs, font=fonts.HEI, size=sz, color=T.PAL.TEXT,
                     anchor=MSO_ANCHOR.MIDDLE)
        y += inch(row_h)
    return top + h


def _place_quote(slide, el, top, stage):
    """引用/重点句：点线边框卡 + 楷体 + 彩色书名号。"""
    sz = T.font_for(stage, "body") + 2
    content = el.get("content", "")
    left = inch(1.2)
    width = inch(13.333 - 2.4)
    n_lines = max(1, -(-len(content) // 26))
    h = inch(0.5 * n_lines + 0.5)
    box = _card_panel(slide, left, top, width, h, fill="FFFDF7",
                      edge_color=T.PAL.ACCENT, edge_w=1.5)
    box.line.dash_style = MSO_LINE.DASH_DOT
    _add_textbox(slide, left + inch(0.3), top, width - inch(0.6), h,
                 runs=[("「", {"size": sz, "color": T.PAL.ACCENT, "bold": True}),
                       * _hl_runs(content, el.get("emphasize"), sz, base_color=T.PAL.TITLE_TEXT),
                       ("」", {"size": sz, "color": T.PAL.ACCENT, "bold": True})],
                 font=fonts.KAI, size=sz, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER, line_spacing=T.L.LINE_SP)
    src = el.get("source")
    if src:
        _add_textbox(slide, left + inch(0.3), top + h - inch(0.05), width - inch(0.6), inch(0.35),
                     f"—— {src}", font=fonts.HEI, size=14, color=T.PAL.TEXT_LIGHT,
                     align=PP_ALIGN.RIGHT)
    return top + h


def _place_ruby_line(slide, el, top, stage):
    """整行注音：Pillow PNG 居中。"""
    png = _render_ruby_png(el.get("text", ""), el.get("ruby", ""), stage,
                           big=bool(el.get("big")))
    if png:
        from PIL import Image as PILImage
        try:
            iw, ih = PILImage.open(io.BytesIO(_png_bytes(png))).size
        except Exception:
            iw, ih = 1200, 200
        max_w = 13.333 - 2 * T.L.MARGIN_X
        w_in = min(max_w, iw / 96)
        h_in = w_in * ih / iw
        x = inch((13.333 - w_in) / 2)
        slide.shapes.add_picture(io.BytesIO(_png_bytes(png)), int(x), int(top),
                                 width=inch(w_in), height=inch(h_in))
        return top + inch(h_in + 0.12)
    _add_textbox(slide, inch(T.L.MARGIN_X), top, inch(13.333 - 2 * T.L.MARGIN_X), inch(0.8),
                 el.get("text", ""), font=fonts.HEI, size=T.font_for(stage, "body"))
    return top + inch(0.9)


def _png_bytes(buf):
    pos = buf.tell()
    buf.seek(0)
    data = buf.read()
    buf.seek(pos)
    return data


def _place_poem(slide, el, top, stage):
    """古诗：白卡 + 居中注音诗句 + 节奏斜线。"""
    stanzas = el.get("stanzas", [])
    title = el.get("title")
    left = inch(1.3)
    width = inch(13.333 - 2.6)
    # 先预渲染各行 PNG 并算出真实高度，再决定卡片总高（避免估算溢出）
    rows = []
    for stanza in stanzas:
        for line in stanza.get("lines", []):
            png = _render_ruby_png(line.get("text", ""), line.get("ruby", ""),
                                   stage, big=True)
            if png:
                from PIL import Image as PILImage
                try:
                    iw, ih = PILImage.open(io.BytesIO(_png_bytes(png))).size
                except Exception:
                    iw, ih = 900, 200
                w_in = min(8.0, iw / 135)
                h_in = w_in * ih / iw
                rows.append((png, w_in, h_in, line.get("text", "")))
            else:
                rows.append((None, 0, 0.7, line.get("text", "")))
    head_h = 1.15 if title else 0.35
    body_h = sum(r[2] + 0.12 for r in rows)
    h_in_total = head_h + body_h + 0.5
    # 空间不足时整体缩放行图（保纵横比），杜绝溢出
    top_in = top / 914400
    avail = T.L.MAX_Y - 0.3 - top_in
    if h_in_total > avail and h_in_total > 0:
        scale = max(0.5, avail / h_in_total)
        rows = [(p, w * scale, hh * scale, t) for (p, w, hh, t) in rows]
        h = inch(head_h * scale + sum(r[2] + 0.12 * scale for r in rows) + 0.4)
    else:
        h = inch(h_in_total)
    _card_panel(slide, left, top, width, h, fill="FFFCF2")
    y = top + inch(0.28)
    if title:
        _add_textbox(slide, left, y, width, inch(0.55),
                     title, font=fonts.KAI, size=T.font_for(stage, "h2") + 2,
                     color=T.PAL.TITLE_TEXT, bold=True, align=PP_ALIGN.CENTER)
        author = el.get("author")
        if author:
            _add_textbox(slide, left, y + inch(0.5), width, inch(0.4),
                         author, font=fonts.KAI, size=16, color=T.PAL.TEXT_LIGHT,
                         align=PP_ALIGN.CENTER)
        y += inch(head_h - 0.15)
    for png, w_in, h_in, text in rows:
        if png:
            x = inch((13.333 - w_in) / 2)
            slide.shapes.add_picture(io.BytesIO(_png_bytes(png)), int(x), int(y),
                                     width=inch(w_in), height=inch(h_in))
        else:
            _add_textbox(slide, left, y, width, inch(0.7),
                         text, font=fonts.KAI,
                         size=T.font_for(stage, "body") + 8, color=T.PAL.TITLE_TEXT,
                         align=PP_ALIGN.CENTER)
        y += inch(h_in + 0.12)
    return top + h


def _place_word_cards(slide, el, top, stage):
    """生字卡：白卡 + 顶部彩带 + 拼音大字 + 汉字特大 + 组词/结构。"""
    cards = el.get("cards", [])
    if not cards:
        return top
    n = len(cards)
    per_row = min(n, 3 if n > 2 else n)
    card_w = (13.333 - 2 * T.L.MARGIN_X - (per_row - 1) * T.L.CARD_GAP) / per_row
    char_sz = T.font_for(stage, "bigchar") - 6
    py_sz = T.font_for(stage, "pinyin") + 4
    n_rows = -(-n // per_row)
    row_h = 2.05
    card_h = inch(row_h)
    for i, card in enumerate(cards):
        r, c = divmod(i, per_row)
        cl = T.inch(T.L.MARGIN_X + c * (card_w + T.L.CARD_GAP))
        ct = top + T.inch(r * (row_h + T.L.CARD_GAP))
        if ct > T.inch(T.L.MAX_Y) - card_h:
            break
        accent = T.PAL.HIGHLIGHTS[i % len(T.PAL.HIGHLIGHTS)]
        _card_panel(slide, cl, ct, T.inch(card_w), card_h,
                    edge_color=T.PAL.BORDER, edge_w=1.0)
        # 顶部彩带（圆角条）
        strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       cl + inch(0.28), ct + inch(0.16),
                                       T.inch(card_w) - inch(0.56), inch(0.09))
        _round_rect(strip, 0.5)
        strip.fill.solid(); strip.fill.fore_color.rgb = _rgb(accent)
        _no_border(strip); strip.shadow.inherit = False
        # 拼音
        pinyin = card.get("pinyin", "")
        tone = py.tone_of(pinyin.split()[0]) if pinyin else 0
        _add_textbox(slide, cl, ct + inch(0.26), T.inch(card_w), inch(0.5),
                     pinyin, font=fonts.MONO, size=py_sz,
                     color=T.PAL.TONE.get(tone, T.PAL.TEXT_LIGHT),
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 大字
        _add_textbox(slide, cl, ct + inch(0.72), T.inch(card_w), inch(char_sz / 72 + 0.25),
                     card.get("char", ""), font=fonts.HEI, size=char_sz,
                     color=T.PAL.TITLE_TEXT, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 组词（彩色）
        groups = "、".join(card.get("groups", [])[:3])
        if groups:
            _add_textbox(slide, cl + inch(0.12), ct + card_h - inch(0.5),
                         T.inch(card_w) - inch(0.24), inch(0.4),
                         groups, font=fonts.HEI, size=15, color=accent, bold=True,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return top + T.inch(n_rows * (row_h + T.L.CARD_GAP))


def _place_revision(slide, el, top, stage):
    """写字指导：田字格 PNG + 拼音 + 结构/部首 + 易错点、运笔要点。"""
    chars = el.get("chars", [])
    if not chars:
        return top
    per_row = 2 if len(chars) > 1 else 1
    cell_w = (13.333 - 2 * T.L.MARGIN_X - (per_row - 1) * T.L.CARD_GAP) / per_row
    n_rows = -(-len(chars) // per_row)
    box_h = 2.35
    for i, ch in enumerate(chars):
        r, c = divmod(i, per_row)
        cl = T.inch(T.L.MARGIN_X + c * (cell_w + T.L.CARD_GAP))
        ct = top + T.inch(r * (box_h + 0.25))
        if ct > T.inch(T.L.MAX_Y) - T.inch(box_h):
            break
        _card_panel(slide, cl, ct, T.inch(cell_w), T.inch(box_h),
                    edge_color=T.PAL.BORDER, edge_w=1.0)
        pinyin = ch.get("pinyin", "")
        tone = py.tone_of(pinyin.split()[0]) if pinyin else 0
        png = _render_tianzi_png(ch.get("char", "？"), stage, tone_color=T.PAL.TITLE_TEXT)
        grid_sz = inch(1.62)
        gx = cl + inch(0.22)
        gy = ct + inch(0.34)
        if png:
            slide.shapes.add_picture(png, int(gx), int(gy), width=grid_sz, height=grid_sz)
        else:
            _add_textbox(slide, int(gx), int(gy), int(grid_sz), int(grid_sz),
                         ch.get("char", ""), font=fonts.KAI, size=60,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        info_x = gx + grid_sz + inch(0.28)
        info_w = cl + T.inch(cell_w) - info_x - inch(0.18)
        # 拼音 + 字
        _add_textbox(slide, int(info_x), ct + inch(0.22), int(info_w), inch(0.42),
                     runs=[(pinyin + "　", {"size": 20, "color": T.PAL.TONE.get(tone, T.PAL.TEXT_LIGHT), "bold": True}),
                           (ch.get("char", ""), {"size": 26, "color": T.PAL.TITLE_TEXT, "bold": True})],
                     font=fonts.HEI, anchor=MSO_ANCHOR.MIDDLE)
        meta_bits = []
        if ch.get("部首"): meta_bits.append("部首 " + str(ch["部首"]))
        if ch.get("结构"): meta_bits.append("结构 " + str(ch["结构"]))
        if ch.get("笔画"): meta_bits.append(str(ch["笔画"]) + "画")
        if meta_bits:
            _add_textbox(slide, int(info_x), ct + inch(0.68), int(info_w), inch(0.32),
                         "｜".join(meta_bits), font=fonts.HEI, size=13,
                         color=T.PAL.TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # 易错 / 运笔（彩色强调）
        y = ct + inch(1.02)
        for key, label, color in (("易错点", "易错", "E2574C"), ("运笔要点", "运笔", "3E8E5A")):
            v = ch.get(key, "")
            if not v:
                continue
            _add_textbox(slide, int(info_x), int(y), int(info_w), inch(0.62),
                         runs=[(f"{label}  ", {"size": 14, "color": color, "bold": True}),
                               (str(v), {"size": 14, "color": T.PAL.TEXT})],
                         font=fonts.HEI, anchor=MSO_ANCHOR.TOP, line_spacing=1.1)
            y += inch(0.62)
    return top + T.inch(n_rows * (box_h + 0.25))


def _place_board(slide, el, top, stage):
    """板书：中心标题 + 树形节点（圆角框 + 连线）。"""
    structure = el.get("structure", [])
    cx = 13.333 / 2
    # 中心标题
    tw = 3.4
    title = el.get("title", "")
    head = _card_panel(slide, inch(cx - tw / 2), top, inch(tw), inch(0.6),
                       fill=T.PAL.ACCENT)
    _add_textbox(slide, inch(cx - tw / 2), top, inch(tw), inch(0.6),
                 title, font=fonts.HEI, size=T.font_for(stage, "h3") + 2,
                 color=T.PAL.TITLE_TEXT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    top += inch(0.85)
    n = len(structure)
    col_w = (13.333 - 2 * 0.9) / max(n, 1)
    box_w = min(col_w - 0.3, 3.3)
    max_bottom = top
    for ni, node in enumerate(structure):
        accent = T.PAL.HIGHLIGHTS[ni % len(T.PAL.HIGHLIGHTS)]
        bx = inch(0.9 + ni * col_w + (col_w - box_w) / 2)
        # 连线：中心底部 → 节点框顶部
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          inch(cx), top - inch(0.25),
                                          int(bx) + inch(box_w / 2), top)
        line.line.color.rgb = _rgb(T.PAL.FOOT)
        line.line.width = Pt(1.4)
        # 一级节点
        nh = 0.52
        _card_panel(slide, int(bx), int(top), inch(box_w), inch(nh),
                    fill="FFFDF7", edge_color=accent, edge_w=1.6)
        _add_textbox(slide, int(bx), int(top), inch(box_w), inch(nh),
                     node.get("node", ""), font=fonts.HEI, size=T.font_for(stage, "body"),
                     color=accent, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y = top + inch(nh + 0.06)
        for child in node.get("children", []):
            ctext = child.get("node", "") if isinstance(child, dict) else str(child)
            ch_lines = max(1, -(-len(ctext) // 11))
            chh = inch(0.34 * ch_lines + 0.16)
            cbx = int(bx) + inch(0.18)
            cbw = inch(box_w - 0.36)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(bx), int(y + chh / 2 - inch(0.05)),
                                         inch(0.1), inch(0.1))
            dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(accent)
            _no_border(dot); dot.shadow.inherit = False
            _add_textbox(slide, cbx, int(y), int(cbw), int(chh),
                         ctext, font=fonts.HEI, size=T.font_for(stage, "body") - 3,
                         color=T.PAL.TEXT, anchor=MSO_ANCHOR.MIDDLE)
            y += chh + inch(0.05)
        max_bottom = max(max_bottom, y)
    return max_bottom + inch(0.1)


def _place_discussion(slide, el, top, stage):
    """讨论卡：绿色提问框 + 提示 + 活动形式徽章。"""
    left = inch(T.L.MARGIN_X)
    width = inch(13.333 - 2 * T.L.MARGIN_X)
    q = el.get("question", "")
    n = max(1, -(-len(q) // 28))
    h = inch(0.5 * n + 1.0)
    _card_panel(slide, left, top, width, h, fill="F3F9F4",
                edge_color=T.PAL.ACCENT2, edge_w=1.4)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + inch(0.28),
                                   top + inch(0.26), inch(0.5), inch(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = _rgb(T.PAL.ACCENT2)
    _no_border(badge); badge.shadow.inherit = False
    _add_textbox(slide, left + inch(0.28), top + inch(0.24), inch(0.5), inch(0.5),
                 "问", font=fonts.HEI, size=18, color=T.PAL.TITLE_TEXT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, left + inch(0.95), top + inch(0.18), width - inch(1.25),
                 inch(0.5 * n + 0.3),
                 q, font=fonts.HEI, size=T.font_for(stage, "h3"),
                 color=T.PAL.TITLE_TEXT, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=T.L.LINE_SP)
    y = top + inch(0.5 * n + 0.6)
    form = el.get("form")
    hint = el.get("hint")
    if form:
        _pill(slide, left + width - inch(0.25), y - inch(0.04), form,
              fill=T.PAL.ACCENT2, color=T.PAL.TITLE_TEXT, size=13)
    if hint:
        _add_textbox(slide, left + inch(0.95), y, width - inch(1.25), inch(0.4),
                     runs=[("💡 ", {"size": 14}),
                           ("提示：" + str(hint), {"size": 14, "color": T.PAL.TEXT_LIGHT})],
                     font=fonts.HEI, color=T.PAL.TEXT_LIGHT)
    return top + h


def _place_table(slide, el, top, stage):
    """表头橙色 + 斑马行，白卡容器。"""
    headers = el.get("headers", [])
    rows = el.get("rows", [])
    if not headers and not rows:
        return top
    n_cols = max(len(headers), max((len(r) for r in rows), default=1), 1)
    width = 13.333 - 2 * T.L.MARGIN_X
    col_w = width / n_cols
    sz = T.font_for(stage, "body") - 2
    row_h = 0.52
    n_rows = len(headers) + len(rows)
    _card_panel(slide, inch(T.L.MARGIN_X), top, inch(width), inch(row_h * n_rows + 0.12))
    y = top + inch(0.06)
    for ci, h in enumerate(headers):
        x = T.inch(T.L.MARGIN_X + ci * col_w)
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), T.inch(col_w), inch(row_h))
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(T.PAL.ACCENT)
        _no_border(cell); cell.shadow.inherit = False
        _add_textbox(slide, int(x), int(y), T.inch(col_w), inch(row_h),
                     str(h), size=sz + 1, color=T.PAL.TITLE_TEXT, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if headers:
        y += inch(row_h)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            if ci >= n_cols:
                break
            x = T.inch(T.L.MARGIN_X + ci * col_w)
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), T.inch(col_w), inch(row_h))
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb("FFFFFF" if ri % 2 == 0 else "FBF3E6")
            cell.line.color.rgb = _rgb(T.PAL.DIVIDER); cell.line.width = Pt(0.75)
            cell.shadow.inherit = False
            _add_textbox(slide, int(x), int(y), T.inch(col_w), inch(row_h),
                         str(val), size=sz, color=T.PAL.TEXT,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += inch(row_h)
    return top + inch(row_h * n_rows + 0.2)


def _place_evaluation(slide, el, top, stage):
    """评价量表：星级行，星星橙黄。"""
    rubric = el.get("rubric", [])
    if not rubric:
        return top
    left = inch(T.L.MARGIN_X)
    width = inch(13.333 - 2 * T.L.MARGIN_X)
    sz = T.font_for(stage, "body") - 2
    row_h = 0.55
    n_rows = len(rubric) + 1
    _card_panel(slide, left, top, width, inch(row_h * n_rows + 0.12))
    y = top + inch(0.06)
    # 表头
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(y), width, inch(row_h))
    head.fill.solid(); head.fill.fore_color.rgb = _rgb(T.PAL.ACCENT2)
    _no_border(head); head.shadow.inherit = False
    _add_textbox(slide, left, int(y), width, inch(row_h), "评价项目　★★★★★",
                 size=sz + 1, color=T.PAL.TITLE_TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y += inch(row_h)
    for i, item in enumerate(rubric):
        crit = item.get("criterion", "")
        levels = item.get("levels", [])
        desc = "；".join(f"{l.get('star','')}★{l.get('desc','')}" if l.get("star") else str(l.get("desc", ""))
                         for l in levels) or ""
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(y), width, inch(row_h))
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb("FFFFFF" if i % 2 == 0 else "F3F8F3")
        _no_border(cell); cell.shadow.inherit = False
        _add_textbox(slide, left + inch(0.25), int(y), inch(4.2), inch(row_h),
                     crit, size=sz, color=T.PAL.TITLE_TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        _add_textbox(slide, left + inch(4.6), int(y), width - inch(4.85), inch(row_h),
                     desc, size=sz - 1, color=T.PAL.TEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += inch(row_h)
    return top + inch(row_h * n_rows + 0.2)


def _place_strokes(slide, el, top, stage):
    """笔顺：每笔一格的田字格序列（Pillow 逐字渲染前 N 笔）。"""
    char = el.get("char", "")
    order = el.get("strokeOrder", [])
    if not char:
        return top
    n = max(len(order), 1)
    cell = 1.35
    total_w = n * cell
    if total_w > 12.0:
        cell = 12.0 / n
        total_w = 12.0
    x0 = (13.333 - total_w) / 2
    h = inch(cell + 0.45)
    png = _render_tianzi_png(char, stage, tone_color=T.PAL.ACCENT)
    x = x0
    for i in range(n):
        lbl = order[i] if i < len(order) else ""
        if png:
            slide.shapes.add_picture(io.BytesIO(_png_bytes(png)), inch(x), top + inch(0.42),
                                     width=inch(cell - 0.18), height=inch(cell - 0.18))
        nb = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x + cell / 2 - 0.17), top,
                                    inch(0.34), inch(0.34))
        nb.fill.solid(); nb.fill.fore_color.rgb = _rgb(T.PAL.ACCENT)
        _no_border(nb); nb.shadow.inherit = False
        _add_textbox(slide, inch(x + cell / 2 - 0.17), top, inch(0.34), inch(0.34),
                     str(i + 1), size=13, color=T.PAL.TITLE_TEXT, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if lbl:
            _add_textbox(slide, inch(x), top + inch(cell + 0.18), inch(cell), inch(0.3),
                         lbl, size=12, color=T.PAL.TEXT_LIGHT, align=PP_ALIGN.CENTER)
        x += cell
    return top + h


# ---- 主渲染 ----
def render(doc: dict, out_path: str) -> str:
    """渲染一份 doc → .pptx。按课时分文件，返回主文件路径。"""
    meta = doc["meta"]
    stage = stage_short(meta["stage"])
    total = meta.get("periods", 1)
    periods = sorted(set(s.get("period", 1) for s in doc["slides"]))

    out_path = Path(out_path)
    main_path = str(out_path)
    paths = []
    for pi, per in enumerate(periods):
        prs = Presentation()
        prs.slide_width = T.SLIDE_W
        prs.slide_height = T.SLIDE_H
        blank = prs.slide_layouts[6]
        slides = [s for s in doc["slides"] if s.get("period", 1) == per]
        n = len(slides)
        for si, sj in enumerate(slides, start=1):
            slide = prs.slides.add_slide(blank)
            layout = classify_slide(sj)
            if layout == "cover":
                if sj.get("kind") == "end":
                    _render_end(slide, sj, meta, stage, si, n)
                else:
                    _render_cover(slide, sj, meta, stage, si, n)
            else:
                _render_content(slide, sj, meta, stage)
                _header_tags(slide, sj, meta, total)
                _footer(slide, meta, si, n)
        fp = main_path if pi == 0 else str(out_path.with_suffix(f".p{per}.pptx"))
        prs.save(fp)
        paths.append(fp)
    return paths[0]


def main(argv=None):
    import json, argparse
    p = argparse.ArgumentParser(description="JSON → PPTX 课件渲染")
    p.add_argument("input", help="课程 JSON 文件")
    p.add_argument("-o", "--output", help="输出 .pptx 路径")
    args = p.parse_args(argv)
    try:
        with open(args.input, encoding="utf-8") as f:
            doc = json.load(f)
        from common.schema import validate
        validate(doc)
    except Exception as e:
        print(f"[pptx] ✗ 输入无效：{e}", file=sys.stderr)
        return 2
    out = args.output or (Path(args.input).with_suffix(".pptx"))
    try:
        main_path = render(doc, out)
        print(f"[pptx] ✓ 已生成：{main_path}")
        return 0
    except Exception as e:
        print(f"[pptx] ✗ 渲染失败：{e}", file=sys.stderr)
        import traceback; traceback.print_exc()
        return 1


if __name__ == "__main__":
    sys.exit(main())
